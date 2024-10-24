from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
import math
import shutil
import os
import formatting_ppt, globalvars

def create_ppt_with_two_images_per_slide(bridgeID, date, image_caption_dict, output_folder, output_pptx_path):
    # Create a new presentation
    prs = Presentation()
    # Set the slide width and height (in inches)
    prs.slide_width = Inches(8.5)
    prs.slide_height = Inches(11)

    # Convert dictionary to a list of tuples for easier iteration
    image_caption_array = list(image_caption_dict.items())

    # Title to the slide position
    title_position = formatting_ppt.Textbox.position(left=1.11, top=0.3, width=6.33, height=0.4)
    # Inspectors position
    inspectors_position = formatting_ppt.Textbox.position(left=0.5, top=0.58, width=7.58, height=0.33)
    # Image positions
    first_img_position = formatting_ppt.Textbox.position(left=1.36, top=0.92, width=5.77, height=4.33)
    first_cap_position = formatting_ppt.Textbox.position(left=0.45, top=5.25, width=7.6, height=0.33)
    second_img_position = formatting_ppt.Textbox.position(left=1.36, top=5.75, width=5.77, height=4.33)
    second_cap_position = formatting_ppt.Textbox.position(left=0.45, top=10.08, width=7.6, height=0.33)
    page_no_position = formatting_ppt.Textbox.position(left=6.09, top=10.42, width=1.98, height=0.36)

    # Loop through the image-caption array in pairs
    for i in range(0, len(image_caption_array), 2):
        # Create a new slide with a blank layout
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        # Remove title placeholder if it exists
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                sp = shape
                sp.element.getparent().remove(sp.element)

        # Add the title
        title_text = f"{bridgeID} Routine Inspection {date[4:6]}/{date[-2:]}/{date[:4]}"
        title_size = 18
        title_textbox = formatting_ppt.Textbox.create_textbox(position=title_position, slide=slide, text=title_text, size=title_size, is_bold=True)

        if i == 0:
            # Add the inspectors caption
            inspectors_text = str(globalvars.inspectors)
            inspectors_size = 12
            inspectors_textbox = formatting_ppt.Textbox.create_textbox(position=inspectors_position, slide=slide, text=inspectors_text, size=inspectors_size, is_bold=False)

        # Image 1
        image_1_name, caption_1 = image_caption_array[i][0], image_caption_array[i][1][0]
        image_1_path = os.path.join(output_folder, image_1_name)
        shape_1 = slide.shapes.add_picture(image_1_path, *first_img_position)
        shape_1.rotation = image_caption_array[i][1][1]
        
        img1_caption_text = str(caption_1) if caption_1 else ""
        img1_caption_size = 12
        img1_caption_textbox = formatting_ppt.Textbox.create_textbox(position=first_cap_position, slide=slide, text=img1_caption_text, size=img1_caption_size, is_bold=False)

        # Check if there is a second image in this pair
        if i + 1 < len(image_caption_array):
            # Image 2
            image_2_name, caption_2 = image_caption_array[i + 1][0], image_caption_array[i + 1][1][0]
            image_2_path = os.path.join(output_folder, image_2_name)
            shape_2 = slide.shapes.add_picture(image_2_path, *second_img_position)
            shape_2.rotation = image_caption_array[i + 1][1][1]  # Ensure we get the correct rotation for the second image
            
            img2_caption_text = str(caption_2) if caption_2 else ""
            img2_caption_size = 12
            img2_caption_textbox = formatting_ppt.Textbox.create_textbox(position=second_cap_position, slide=slide, text=img2_caption_text, size=img2_caption_size, is_bold=False)

        # Add the page numbering
        page_no_text = f"{int((i + 2) / 2)}/{math.ceil(len(image_caption_array) / 2)}"
        page_no_size = 18
        page_no_textbox = formatting_ppt.Textbox.create_textbox(position=page_no_position, slide=slide, text=page_no_text, size=page_no_size, is_bold=True)

    # Save the new PowerPoint presentation
    prs.save(output_pptx_path)

    # Ensure the output folder is deleted
    if os.path.exists(output_folder):
        shutil.rmtree(output_folder, ignore_errors=False)
